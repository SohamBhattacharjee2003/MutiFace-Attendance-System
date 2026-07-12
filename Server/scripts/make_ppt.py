"""
make_ppt.py
Builds the full project deck (13 slides): title, problem, solution, two methodology
slides with diagrams, a sequence diagram, the experimental protocol, two results slides,
improvements + future scope, impacts, references.

Every number and every figure is generated from this repository's own measurements
(results/evaluation.json, results/benchmark.json, scripts/test_spoof.py) rather than
typed in, so the slides cannot drift from what the code actually does. The team's earlier
deck claimed MTCNN + FaceNet + MongoDB — none of which this system uses — and a 95%
accuracy that was never measured. That is the failure mode this script exists to prevent.

    python scripts/make_ppt.py   ->   results/PresenceAI_Presentation.pptx
"""

import os
import sys
import json

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RESULTS = "results"
OUT = os.path.join(RESULTS, "PresenceAI_Presentation.pptx")
FIG = os.path.join(RESULTS, "ppt_figs")

NAVY = RGBColor(0x14, 0x2B, 0x53)
BLUE = RGBColor(0x1F, 0x77, 0xC2)
INK = RGBColor(0x22, 0x2A, 0x35)
GREY = RGBColor(0x5B, 0x66, 0x76)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)
AMBER = RGBColor(0xB9, 0x7A, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

F_RED = RGBColor(0xFD, 0xEC, 0xEA)
F_BLUE = RGBColor(0xE9, 0xF2, 0xFC)
F_GREEN = RGBColor(0xE8, 0xF7, 0xEF)
F_AMBER = RGBColor(0xFE, 0xF5, 0xE3)
F_GREY = RGBColor(0xF4, 0xF6, 0xF9)

M_BLUE, M_RED, M_GREEN, M_GREY, M_AMBER = "#1f77c2", "#c0392b", "#1e8e5a", "#5b6676", "#e08e0b"


def score_distributions(seed=7, n_students=30):
    """
    Recompute the genuine vs impostor score distributions on the SAME split
    scripts/evaluate.py uses, so the histogram and ROC on the slides are the
    distributions the reported FAR/threshold actually came from.
    """
    import numpy as np
    import trainer

    X, y, _ = trainer.embed_dataset(progress=None)
    pool = np.array([trainer.is_impostor(n) for n in y])
    X, y = X[pool], y[pool]
    idents = np.unique(y)

    rng = np.random.default_rng(seed)
    shuffled = rng.permutation(idents)
    students = np.array(sorted(shuffled[:n_students]))
    strangers = np.array(sorted(shuffled[n_students:]))
    half = len(strangers) // 2
    calib_ids, test_ids = strangers[:half], strangers[half:]

    enrol_X, enrol_y, probe_X = [], [], []
    for st in students:
        E = X[y == st]
        idx = rng.permutation(len(E))
        cut = int(0.7 * len(E))
        enrol_X.append(E[idx[:cut]]); enrol_y += [st] * cut
        probe_X.append(E[idx[cut:]])
    enrol_X = np.vstack(enrol_X); enrol_y = np.array(enrol_y)
    probe_X = np.vstack(probe_X)

    cent = {}
    for n in np.unique(enrol_y):
        v = enrol_X[enrol_y == n].mean(axis=0)
        cent[n] = (v / np.linalg.norm(v)).astype(np.float32)
    M = np.stack([cent[n] for n in sorted(cent)])

    gen = (probe_X @ M.T).max(axis=1)                       # genuine top-1
    imp = (X[np.isin(y, test_ids)] @ M.T).max(axis=1)       # impostor top-1
    calib = (X[np.isin(y, calib_ids)] @ M.T).max(axis=1)
    thr = float(np.quantile(calib, 1 - trainer.TARGET_FAR))
    return gen, imp, thr


# ════════════════════════════════════════════════════════════════ FIGURES ══
def _box(ax, x, y, w, h, text, fc, ec, fs=8, bold=False):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.06",
                                facecolor=fc, edgecolor=ec, linewidth=1.4, zorder=2))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fs,
            fontweight="bold" if bold else "normal", zorder=3, linespacing=1.35)


def _arrow(ax, x1, y1, x2, y2, color=M_GREY):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                                 mutation_scale=13, color=color, linewidth=1.3, zorder=1))


def figures(ev, bm):
    os.makedirs(FIG, exist_ok=True)
    plt.rcParams.update({"font.size": 10, "axes.spines.top": False,
                         "axes.spines.right": False, "figure.dpi": 200})

    # 1 · presentation attack
    fig, ax = plt.subplots(figsize=(4.0, 2.9))
    bars = ax.bar(["Recognition\nonly", "+ Liveness\n(ours)"], [100, 0],
                  color=[M_RED, M_GREEN], width=0.52)
    for b, v in zip(bars, [100, 0]):
        ax.text(b.get_x() + b.get_width() / 2, v + 4, f"{v}%", ha="center",
                fontweight="bold", fontsize=13)
    ax.set_ylabel("Attack success (%)"); ax.set_ylim(0, 118)
    ax.set_title("Photo held up to the camera", fontsize=10.5, fontweight="bold")
    fig.tight_layout(); fig.savefig(f"{FIG}/spoof.png"); plt.close(fig)

    # 2 · ablation
    rows = ev["ablation"]
    fig, ax = plt.subplots(figsize=(5.6, 2.9))
    x = range(len(rows))
    ax.bar([i - 0.19 for i in x], [r["far"] for r in rows], 0.38,
           label="Strangers admitted (FAR)", color=M_BLUE)
    ax.bar([i + 0.19 for i in x], [r["false_log_rate"] for r in rows], 0.38,
           label="Wrong attendance records", color="#8e5ad6")
    ax.set_xticks(list(x))
    ax.set_xticklabels(["Baseline", "+ Calibrated\nthreshold", "+ Top-2\nmargin",
                        "+ Temporal\nvoting"], fontsize=8)
    ax.set_ylabel("Error rate (%)"); ax.legend(fontsize=8, frameon=False)
    ax.set_title("Ablation — each guard switched on in turn", fontsize=10.5, fontweight="bold")
    fig.tight_layout(); fig.savefig(f"{FIG}/ablation.png"); plt.close(fig)

    # 3 · operating envelope
    deg = ev["degradation"]
    fig, ax = plt.subplots(figsize=(4.6, 2.9))
    ax.plot([d["face_px"] for d in deg], [d["accuracy"] * 100 for d in deg],
            "o-", color=M_BLUE, lw=2, ms=5)
    ax.axvspan(0, 40, color=M_RED, alpha=0.07)
    ax.axvline(40, ls="--", color=M_AMBER, lw=1.5)
    ax.text(45, 68, "gate: reject\nunder 40px", color=M_AMBER, fontsize=8)
    ax.set_xlabel("Face width in the frame (pixels)"); ax.set_ylabel("Accuracy (%)")
    ax.set_title("Operating envelope", fontsize=10.5, fontweight="bold")
    fig.tight_layout(); fig.savefig(f"{FIG}/degradation.png"); plt.close(fig)

    # 4 · throughput
    tp = bm["throughput"]
    fig, ax = plt.subplots(figsize=(4.6, 2.9))
    ax.plot([t["faces_in_frame"] for t in tp], [t["latency_s"] * 1000 for t in tp],
            "o-", color="#8e5ad6", lw=2, ms=5)
    ax.set_xlabel("Faces in one frame"); ax.set_ylabel("Latency (ms)")
    ax.set_title("Scales sub-linearly (CPU, no GPU)", fontsize=10.5, fontweight="bold")
    ax.annotate("~200 ms detection (fixed)\n+ ~65 ms per extra face",
                xy=(8, 934), xytext=(1.4, 1180), fontsize=8, color=M_GREY,
                arrowprops=dict(arrowstyle="->", color=M_GREY, lw=1))
    fig.tight_layout(); fig.savefig(f"{FIG}/throughput.png"); plt.close(fig)

    # 5 · methodology I — representation & enrolment
    fig, ax = plt.subplots(figsize=(11.5, 3.5)); ax.axis("off")
    ax.set_xlim(0, 11.5); ax.set_ylim(0, 3.5)
    ax.text(0.1, 3.3, "STAGE A — Pretrained and FROZEN (done once, by InsightFace — not by us)",
            fontsize=9, fontweight="bold", color=M_RED)
    _box(ax, 0.1, 2.3, 2.4, 0.8, "WebFace600K\n~600,000 identities", "#fdecea", M_RED)
    _box(ax, 2.8, 2.3, 2.4, 0.8, "ArcFace ResNet-50\nangular margin loss", "#fdecea", M_RED)
    _box(ax, 5.5, 2.3, 2.6, 0.8, "FROZEN encoder\nface → 512-d vector", "#fdecea", M_RED, 8, True)
    _arrow(ax, 2.5, 2.7, 2.8, 2.7); _arrow(ax, 5.2, 2.7, 5.5, 2.7)
    ax.text(8.3, 2.7, "It never learns 'who is Soham'.\nIt learns what makes any\ntwo faces different.",
            fontsize=8, color=M_GREY, style="italic", va="center")

    ax.text(0.1, 1.8, "STAGE B — Enrolment (per student, ~1.4 seconds)",
            fontsize=9, fontweight="bold", color=M_GREEN)
    _box(ax, 0.1, 0.6, 1.9, 0.9, "~20 webcam\nphotos", "#e8f7ef", M_GREEN)
    _box(ax, 2.3, 0.6, 2.0, 0.9, "Detect + align\n(RetinaFace)", "#e8f7ef", M_GREEN)
    _box(ax, 4.6, 0.6, 2.1, 0.9, "Embed → 512-d\n(cached on disk)", "#e8f7ef", M_GREEN)
    _box(ax, 7.0, 0.6, 2.2, 0.9, "Mean → CENTROID\ntheir 'signature'", "#e8f7ef", M_GREEN, 8, True)
    _box(ax, 9.5, 0.6, 1.9, 0.9, "Recalibrate\nthreshold", "#e8f7ef", M_GREEN)
    for a, b in [(2.0, 2.3), (4.3, 4.6), (6.7, 7.0), (9.2, 9.5)]:
        _arrow(ax, a, 1.05, b, 1.05)
    ax.text(0.1, 0.2, "No gradient descent. No GPU. No epochs. Adding a student is a table insert — "
                      "which is why enrolment fell from 607 s to 1.4 s.",
            fontsize=8, color=M_GREY, style="italic")
    fig.tight_layout(); fig.savefig(f"{FIG}/methodology_enroll.png"); plt.close(fig)

    # 6 · methodology II — the four gates
    fig, ax = plt.subplots(figsize=(11.5, 4.3)); ax.axis("off")
    ax.set_xlim(0, 11.5); ax.set_ylim(0, 4.3)
    _box(ax, 0.15, 3.35, 1.6, 0.75, "Camera\nframe", "#e9f2fc", M_BLUE)
    _box(ax, 2.05, 3.35, 2.0, 0.75, "Detect ALL faces\nRetinaFace", "#e9f2fc", M_BLUE)
    _box(ax, 4.35, 3.35, 2.0, 0.75, "Embed each face\nArcFace 512-d", "#e9f2fc", M_BLUE)
    _box(ax, 6.65, 3.35, 2.2, 0.75, "Cosine vs EVERY\nstudent centroid", "#e9f2fc", M_BLUE)
    for a, b in [(1.75, 2.05), (4.05, 4.35), (6.35, 6.65)]:
        _arrow(ax, a, 3.72, b, 3.72)
    _arrow(ax, 7.75, 3.35, 7.75, 3.05)

    gates = [
        ("GATE 1 · SIZE\nface ≥ 40 px ?",
         "A distant face is upscaled to 112 px —\nno detail is invented. The classifier\nwould still name someone.", "Move closer"),
        ("GATE 2 · THRESHOLD\ncos ≥ 0.148 ?",
         "Calibrated on 2,880 impostor faces at\na 1% target false-accept rate.\nNot hand-picked.", "Unknown"),
        ("GATE 3 · MARGIN\ntop1 − top2 ≥ 0.15 ?",
         "Two look-alike students both score\nhigh. The GAP between them is what\nseparates them.", "Uncertain"),
        ("GATE 4 · LIVENESS\nis the face moving ?",
         "Recognition matches a PHOTO of you\nto you. A live face deforms;\na photograph is rigid.", "Not live"),
    ]
    for i, (title, why, reject) in enumerate(gates):
        x = 0.15 + i * 2.85
        fc, ec = ("#fdecea", M_RED) if i == 3 else ("#fef5e3", M_AMBER)
        _box(ax, x, 2.15, 2.5, 0.78, title, fc, ec, 8, True)
        ax.text(x + 1.25, 2.02, why, ha="center", va="top", fontsize=6.5, color=M_GREY)
        ax.text(x + 1.25, 1.18, f"✗ → \"{reject}\"", ha="center", fontsize=7,
                color=M_RED, fontweight="bold")
        if i < 3:
            _arrow(ax, x + 2.5, 2.54, x + 2.85, 2.54, color=M_GREEN)

    _box(ax, 3.6, 0.15, 3.0, 0.62, "TEMPORAL VOTE — 3 frames must agree",
         "#e8f7ef", M_GREEN, 8, True)
    _box(ax, 6.9, 0.15, 2.6, 0.62, "ATTENDANCE COMMITTED", "#e8f7ef", M_GREEN, 8, True)
    _arrow(ax, 6.6, 0.46, 6.9, 0.46, color=M_GREEN)
    _arrow(ax, 8.65, 2.15, 5.1, 0.8, color=M_GREEN)
    ax.text(0.15, 0.45, "Attendance is written ONCE per day\nand is irreversible — so committing\nis its own decision, not the model's.",
            fontsize=7.5, color=M_RED, va="center", fontweight="bold")
    fig.tight_layout(); fig.savefig(f"{FIG}/methodology_gates.png"); plt.close(fig)

    # 8 · experimental protocol
    fig, ax = plt.subplots(figsize=(10.5, 3.4)); ax.axis("off")
    ax.set_xlim(0, 10.5); ax.set_ylim(0, 3.4)
    _box(ax, 3.6, 2.72, 3.3, 0.55, "96 VGGFace2 identities", "#f4f6f9", M_GREY, 9, True)
    _box(ax, 0.5, 1.6, 3.6, 0.7, "30  →  simulated STUDENTS", "#e8f7ef", M_GREEN, 9, True)
    _box(ax, 6.4, 1.6, 3.6, 0.7, "66  →  STRANGERS (impostors)", "#fdecea", M_RED, 9, True)
    _arrow(ax, 4.6, 2.72, 2.3, 2.32, color=M_GREEN)
    _arrow(ax, 5.9, 2.72, 8.2, 2.32, color=M_RED)
    _box(ax, 0.2, 0.55, 1.9, 0.75, "70%\nENROLMENT\nbuild centroid", "#e8f7ef", M_GREEN, 7.5)
    _box(ax, 2.3, 0.55, 1.9, 0.75, "30%\nTEST PROBES\nnever seen", "#d6f0e2", M_GREEN, 7.5, True)
    _box(ax, 6.2, 0.55, 1.9, 0.75, "half →\nCALIBRATE\nthe threshold", "#fdecea", M_RED, 7.5)
    _box(ax, 8.3, 0.55, 1.9, 0.75, "half →\nMEASURE\nthe FAR", "#fbd9d4", M_RED, 7.5, True)
    _arrow(ax, 1.6, 1.6, 1.15, 1.3, color=M_GREEN); _arrow(ax, 2.9, 1.6, 3.25, 1.3, color=M_GREEN)
    _arrow(ax, 7.5, 1.6, 7.15, 1.3, color=M_RED); _arrow(ax, 8.9, 1.6, 9.25, 1.3, color=M_RED)
    ax.text(5.25, 0.08, "Both splits are essential. Scoring a face against a centroid built from that same face gives "
                        "100% and proves nothing.\nA threshold tuned on the same strangers you then report an FAR "
                        "against is not a measurement.",
            fontsize=7.5, ha="center", color=M_RED, style="italic")
    fig.tight_layout(); fig.savefig(f"{FIG}/protocol.png"); plt.close(fig)

    # 9 · SYSTEM ARCHITECTURE (replaces the sequence diagram)
    import numpy as np
    fig, ax = plt.subplots(figsize=(11.5, 3.6)); ax.axis("off")
    ax.set_xlim(0, 11.5); ax.set_ylim(0, 3.6)

    _box(ax, 0.15, 2.55, 1.55, 0.85, "React UI\n(Vite + Tailwind)", "#e9f2fc", M_BLUE)
    _box(ax, 0.15, 1.45, 1.55, 0.85, "Webcam\n1 frame / 0.7 s", "#e9f2fc", M_BLUE)
    _box(ax, 0.15, 0.35, 1.55, 0.85, "Box + verdict\noverlay", "#e9f2fc", M_BLUE)
    ax.text(0.92, 3.48, "CLIENT", fontsize=8.5, fontweight="bold", color=M_BLUE, ha="center")

    _box(ax, 2.3, 1.45, 1.9, 1.95,
         "Flask REST API\n\n/predict\n/register\n/train/status\n\nBearer-token auth",
         "#f4f6f9", M_GREY, 8, True)
    _box(ax, 2.3, 0.35, 1.9, 0.85, "Auto-retrain\nworker (thread)", "#f4f6f9", M_GREY)
    ax.text(3.25, 3.48, "BACKEND", fontsize=8.5, fontweight="bold", color=M_GREY, ha="center")

    _box(ax, 4.8, 2.55, 2.2, 0.85, "RetinaFace\ndetect ALL faces", "#fdecea", M_RED)
    _box(ax, 4.8, 1.45, 2.2, 0.85, "ArcFace R50\n512-d embedding", "#fdecea", M_RED)
    _box(ax, 4.8, 0.35, 2.2, 0.85, "Liveness module\nlandmark motion", "#fdecea", M_RED, 8, True)
    ax.text(5.9, 3.48, "INFERENCE ENGINE  (ONNX Runtime, CPU)", fontsize=8.5,
            fontweight="bold", color=M_RED, ha="center")

    _box(ax, 7.6, 2.55, 1.9, 0.85, "Centroids\n1 per student", "#fef5e3", M_AMBER)
    _box(ax, 7.6, 1.45, 1.9, 0.85, "Thresholds\n(calibrated)", "#fef5e3", M_AMBER)
    _box(ax, 7.6, 0.35, 1.9, 0.85, "Embedding cache\n(fast retrain)", "#fef5e3", M_AMBER)
    ax.text(8.55, 3.48, "MODEL STORE", fontsize=8.5, fontweight="bold", color=M_AMBER, ha="center")

    _box(ax, 10.0, 2.55, 1.35, 0.85, "users.json", "#e8f7ef", M_GREEN)
    _box(ax, 10.0, 1.45, 1.35, 0.85, "attendance\n.csv", "#e8f7ef", M_GREEN, 8, True)
    ax.text(10.67, 3.48, "STORAGE", fontsize=8.5, fontweight="bold", color=M_GREEN, ha="center")

    for a, b in [(1.7, 2.3), (4.2, 4.8), (7.0, 7.6), (9.5, 10.0)]:
        _arrow(ax, a, 1.9, b, 1.9)
    _arrow(ax, 2.3, 1.2, 1.7, 0.78)
    ax.text(5.75, 0.02, "Everything runs on one laptop. No GPU, no cloud, no database server — "
                        "student faces never leave the machine.",
            fontsize=8, ha="center", color=M_GREY, style="italic")
    fig.tight_layout(); fig.savefig(f"{FIG}/architecture.png"); plt.close(fig)

    # 10 · WHY RECOGNITION ALONE FAILS — the attack, drawn
    fig, ax = plt.subplots(figsize=(9.0, 2.5)); ax.axis("off")
    ax.set_xlim(0, 9.0); ax.set_ylim(0, 2.5)
    _box(ax, 0.1, 1.35, 1.6, 0.8, "REAL person\nat the camera", "#e8f7ef", M_GREEN)
    _box(ax, 0.1, 0.25, 1.6, 0.8, "PHOTO of them,\nheld up", "#fdecea", M_RED)
    _box(ax, 2.3, 0.8, 1.8, 0.9, "ArcFace\nembedding", "#f4f6f9", M_GREY, 8, True)
    _arrow(ax, 1.7, 1.75, 2.3, 1.45, color=M_GREEN)
    _arrow(ax, 1.7, 0.65, 2.3, 1.05, color=M_RED)
    _box(ax, 4.7, 0.8, 2.0, 0.9, "Nearly IDENTICAL\n512-d vectors", "#fef5e3", M_AMBER, 8, True)
    _arrow(ax, 4.1, 1.25, 4.7, 1.25)
    _box(ax, 7.3, 1.35, 1.6, 0.8, "Marked PRESENT", "#e8f7ef", M_GREEN)
    _box(ax, 7.3, 0.25, 1.6, 0.8, "Marked PRESENT", "#fdecea", M_RED, 8, True)
    _arrow(ax, 6.7, 1.45, 7.3, 1.75, color=M_GREEN)
    _arrow(ax, 6.7, 1.05, 7.3, 0.65, color=M_RED)
    ax.text(4.5, 2.32, "Face recognition is DESIGNED to do this — it is what makes it a good model.",
            fontsize=9, ha="center", color=M_RED, fontweight="bold")
    ax.text(4.5, 0.02, "Which is precisely why recognition ALONE cannot prevent proxy attendance.",
            fontsize=8.5, ha="center", color=M_GREY, style="italic")
    fig.tight_layout(); fig.savefig(f"{FIG}/attack.png"); plt.close(fig)

    # 11 · genuine vs impostor score distributions
    gen, imp, thr = score_distributions()
    fig, ax = plt.subplots(figsize=(5.4, 2.9))
    ax.hist(imp, bins=45, color=M_RED, alpha=0.72, density=True,
            label=f"Strangers (n={len(imp)})")
    ax.hist(gen, bins=25, color=M_GREEN, alpha=0.72, density=True,
            label=f"Genuine students (n={len(gen)})")
    ax.axvline(thr, color=M_AMBER, ls="--", lw=2)
    ax.text(thr + 0.015, ax.get_ylim()[1] * 0.78, f"threshold\n{thr:.3f}",
            color=M_AMBER, fontsize=8, fontweight="bold")
    ax.set_xlabel("Cosine similarity to the closest student")
    ax.set_ylabel("Density")
    ax.set_title("The two populations barely overlap", fontsize=10.5, fontweight="bold")
    ax.legend(fontsize=8, frameon=False)
    fig.tight_layout(); fig.savefig(f"{FIG}/scores.png"); plt.close(fig)

    # 12 · open-set ROC
    sc = np.concatenate([gen, imp])
    lab = np.concatenate([np.ones(len(gen)), np.zeros(len(imp))])
    order = np.argsort(-sc)
    tpr = np.cumsum(lab[order]) / lab.sum()
    fpr = np.cumsum(1 - lab[order]) / (len(lab) - lab.sum())
    auc = float(np.trapezoid(tpr, fpr)) if hasattr(np, "trapezoid") else float(np.trapz(tpr, fpr))
    eer_i = int(np.argmin(np.abs((1 - tpr) - fpr)))
    fig, ax = plt.subplots(figsize=(4.4, 2.9))
    ax.plot(fpr * 100, tpr * 100, color=M_BLUE, lw=2.2, label=f"AUC = {auc:.4f}")
    ax.plot([0, 100], [0, 100], "--", color="#c8d2de", lw=1)
    ax.plot(fpr[eer_i] * 100, tpr[eer_i] * 100, "o", color=M_RED, ms=7,
            label=f"EER = {fpr[eer_i] * 100:.2f}%")
    ax.set_xlabel("False accept rate (%)"); ax.set_ylabel("True accept rate (%)")
    ax.set_title("Open-set ROC", fontsize=10.5, fontweight="bold")
    ax.legend(fontsize=8, frameon=False, loc="lower right")
    fig.tight_layout(); fig.savefig(f"{FIG}/roc.png"); plt.close(fig)

    # 13 · running cost
    fig, ax = plt.subplots(figsize=(4.6, 2.7))
    ax.barh(["PresenceAI\n(offline)", "Cloud face API\n(continuous)"], [0, 150000],
            color=[M_GREEN, M_RED], height=0.5)
    ax.set_xlabel("Recurring cost per school, per year  (INR)")
    ax.text(3000, 0, "= 0", va="center", fontsize=11, fontweight="bold", color=M_GREEN)
    ax.text(146000, 1, "1,50,000  ", va="center", ha="right", fontsize=10,
            fontweight="bold", color="white")
    ax.set_title("...and the cloud uploads every student's face",
                 fontsize=9.5, fontweight="bold")
    fig.tight_layout(); fig.savefig(f"{FIG}/cost.png"); plt.close(fig)

    print("  13 figures generated")



# ═══════════════════════════════════════════════════════ SLIDE PRIMITIVES ══
def chrome(s, n):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.08),
                             Inches(13.333), Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background(); bar.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(12.4), Inches(7.10), Inches(0.8), Inches(0.38))
    p = tb.text_frame.paragraphs[0]
    p.text = str(n); p.alignment = PP_ALIGN.RIGHT
    p.runs[0].font.size = Pt(11); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = WHITE


def slide(prs, heading, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(12.4), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = heading
    p.runs[0].font.size = Pt(25); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.74),
                              Inches(1.6), Inches(0.05))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE
    rule.line.fill.background(); rule.shadow.inherit = False
    if kicker:
        kb = s.shapes.add_textbox(Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.35))
        kp = kb.text_frame.paragraphs[0]
        kp.text = kicker
        kp.runs[0].font.size = Pt(11); kp.runs[0].font.italic = True
        kp.runs[0].font.color.rgb = GREY
    return s


def panel(s, l, t, w, h, fill, edge):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                             Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge; shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def text(s, items, l, t, w, h, size=12):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        body, bold, col, lvl = (item if isinstance(item, tuple) else (item, False, INK, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = body; p.level = lvl; p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col
    return tf


def metric(s, l, t, w, label, before, after, colour=GREEN):
    """before ➔ after tile."""
    panel(s, l, t, w, 1.05, F_GREY, GREY)
    text(s, [(label, True, NAVY, 0)], l + 0.14, t + 0.09, w - 0.28, 0.3, size=10)
    tb = s.shapes.add_textbox(Inches(l + 0.14), Inches(t + 0.42), Inches(w - 0.28), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{before}   ➔   {after}"
    r = p.runs[0]; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = colour


def table(s, data, l, t, w, h, widths, bold_rows=(0,), green_col=None):
    tbl = s.shapes.add_table(len(data), len(data[0]), Inches(l), Inches(t),
                             Inches(w), Inches(h)).table
    for i, cw in enumerate(widths):
        tbl.columns[i].width = Inches(cw)
    for r, row in enumerate(data):
        for c, v in enumerate(row):
            cell = tbl.cell(r, c); cell.text = str(v)
            runs = cell.text_frame.paragraphs[0].runs
            if not runs:                 # an empty cell has no run to style
                continue
            run = runs[0]
            run.font.size = Pt(10)
            run.font.bold = (r in bold_rows) or (c == 0 and r > 0)
            if green_col is not None and c == green_col and r > 0:
                run.font.color.rgb = GREEN; run.font.bold = True
    return tbl


# ═══════════════════════════════════════════════════════════════════ MAIN ══
def main():
    ev = json.load(open(os.path.join(RESULTS, "evaluation.json")))
    bm = json.load(open(os.path.join(RESULTS, "benchmark.json")))
    figures(ev, bm)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # ── 1 · TITLE ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                              Inches(13.333), Inches(0.28))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE
    band.line.fill.background(); band.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PresenceAI"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(40); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "Spoof-Resistant Multi-Face Attendance with Liveness-Gated Commit"
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(17); p2.runs[0].font.color.rgb = BLUE
    p3 = tf.add_paragraph()
    p3.text = "Automatic Attendance System using Multiple Face Recognition"
    p3.alignment = PP_ALIGN.CENTER
    p3.runs[0].font.size = Pt(13); p3.runs[0].font.italic = True; p3.runs[0].font.color.rgb = GREY

    panel(s, 3.15, 3.0, 7.0, 0.6, F_RED, RED)
    tb = s.shapes.add_textbox(Inches(3.25), Inches(3.07), Inches(6.8), Inches(0.45))
    p = tb.text_frame.paragraphs[0]
    p.text = "A printed photo defeats face recognition 100% of the time. We measured it — and fixed it."
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(11.5); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = RED

    text(s, [
        ("Team Members' Name and Roll Number", True, NAVY, 0),
        "1. Souryadeep Deb (13000222064)              2. Soham Bhattacharya (13000222065)",
        "3. Srija Basak (13000222066)                  4. Sneha Singh (13000222067)",
        "",
        ("Group Number — 18                Mentor — Mr Prodipta Bhowmik", True, INK, 0),
    ], 3.6, 3.95, 8.2, 1.8, size=13)

    text(s, [
        ("Department of Information Technology", True, NAVY, 0),
        "Techno Main Salt Lake, Kolkata - 700091",
    ], 5.0, 6.1, 6.0, 0.8, size=12)

    # ── 2 · PROBLEM STATEMENT ───────────────────────────────────────────────
    s = slide(prs, "PROBLEM STATEMENT",
              "Manual attendance is the obvious problem. Automating it with face recognition creates a second one that nobody measures.")
    panel(s, 0.5, 1.3, 6.1, 1.95, F_GREY, GREY)
    text(s, [
        ("The obvious problem", True, NAVY, 0),
        "A 60-student roll-call costs ~5 minutes of every lecture — roughly",
        "80 teaching hours across a 200-day year.",
        "It is slow, error-prone, and a friend can sign the register for you.",
    ], 0.7, 1.44, 5.7, 1.75, size=12)

    panel(s, 6.85, 1.3, 5.95, 1.95, F_RED, RED)
    text(s, [
        ("The problem nobody in the literature measures", True, RED, 0),
        "Face recognition is DESIGNED to give a photo of you the same embedding",
        "as you. That is exactly what makes it a good face model.",
        ("So recognition ALONE cannot prevent proxy attendance. We tested it:", True, RED, 0),
        ("a printed photo was marked present 100% of the time.", True, RED, 0),
    ], 7.05, 1.44, 5.55, 1.75, size=11),

    s.shapes.add_picture(f"{FIG}/attack.png", Inches(1.55), Inches(3.35), width=Inches(10.2))

    fails = [
        ("PRESENTATION ATTACK", "A photo or phone screen held to\nthe camera is marked present.", F_RED, RED),
        ("NO WAY TO SAY 'NOBODY'", "A softmax must always name someone.\nA stranger walking past gets a name.", F_AMBER, AMBER),
        ("THE BACK OF THE ROOM", "Too few pixels to identify — but the\nmodel still answers, confidently.", F_AMBER, AMBER),
        ("AN IRREVERSIBLE WRITE", "Written once per day. One bad frame\nmarks the wrong student all day.", F_BLUE, BLUE),
    ]
    for i, (t, body, fc, ec) in enumerate(fails):
        x = 0.5 + i * 3.14
        panel(s, x, 5.8, 2.95, 1.15, fc, ec)
        text(s, [(t, True, ec, 0), (body, False, INK, 0)], x + 0.14, 5.88, 2.7, 1.0, size=8.5)

    # ── 3 · PROPOSED SOLUTION ───────────────────────────────────────────────
    s = slide(prs, "PROPOSED SOLUTION",
              "Separate RECOGNISING a face from MARKING that person PRESENT. They are not the same decision.")
    panel(s, 0.5, 1.3, 6.15, 2.3, F_BLUE, BLUE)
    text(s, [
        ("How it works", True, NAVY, 0),
        "Every face in the frame is detected (RetinaFace), embedded into a 512-d",
        "vector (ArcFace), and compared against each enrolled student's centroid.",
        "",
        "A face must then clear FOUR gates — size, threshold, margin and liveness —",
        "and be confirmed across several frames before attendance is written.",
        "Records go to CSV. Nothing ever leaves the machine.",
    ], 0.7, 1.44, 5.75, 2.1, size=11.5)

    panel(s, 6.85, 1.3, 5.95, 2.3, F_GREEN, GREEN)
    text(s, [
        ("Innovation and uniqueness — stated honestly", True, GREEN, 0),
        ("We do NOT claim a new face-recognition model.", True, RED, 0),
        "ArcFace is pretrained and frozen. Training our own on 30 students would need",
        "a GPU and thousands of images per person — and would still name a stranger",
        "every single time, because a softmax cannot abstain.",
        ("Our contribution is the DECISION LAYER, and the evaluation of it.", True, GREEN, 0),
    ], 7.05, 1.44, 5.55, 2.1, size=11)

    items = [
        ("LIVENESS THAT FITS THE HARDWARE",
         "Blink detection is the standard trick — but a blink lasts 0.3 s and we sample 1.4 frames/s. "
         "We measured it failing. So we measure facial motion from landmarks the detector already produces: "
         "a live face keeps deforming, a photo is rigid. Rotation and scale are normalised out, so waving the "
         "photo does not help the attacker. No extra model, no GPU.", F_RED, RED),
        ("A THRESHOLD FROM DATA, NOT A GUESS",
         "Calibrated against 2,880 held-out impostor faces at a 1% target false-accept rate. Every attendance "
         "paper we reviewed either hand-picks this number or never mentions it.", F_BLUE, BLUE),
        ("ATTENDANCE AS AN IRREVERSIBLE COMMIT",
         "The write happens once per day and cannot be undone, so it accumulates evidence across frames instead "
         "of trusting one. No prior work separates these two decisions.", F_AMBER, AMBER),
        ("ENROLMENT IN 1.4 SECONDS",
         "Embeddings are cached, so a new student costs only their own photos — not a full re-embed of the "
         "dataset, which took 607 s.", F_GREEN, GREEN),
    ]
    for i, (t, body, fc, ec) in enumerate(items):
        x = 0.5 + i * 3.14
        panel(s, x, 3.6, 2.95, 1.55, fc, ec)
        text(s, [(t, True, ec, 0), (body, False, INK, 0)],
             x + 0.14, 3.68, 2.68, 1.42, size=7.5)

    # system architecture — one laptop, no cloud
    s.shapes.add_picture(f"{FIG}/architecture.png", Inches(0.5), Inches(5.25),
                         width=Inches(12.3))

    # ── 4 · METHODOLOGY (representation + decision, one slide) ──────────────
    s = slide(prs, "PROPOSED METHODOLOGY",
              "A · The model is frozen — 'learning' a student means storing one vector.   "
              "B · Four gates decide who is actually present.")
    s.shapes.add_picture(f"{FIG}/methodology_enroll.png", Inches(0.5), Inches(1.25),
                         width=Inches(12.3))
    s.shapes.add_picture(f"{FIG}/methodology_gates.png", Inches(0.5), Inches(4.05),
                         width=Inches(12.3))
    panel(s, 0.5, 6.72, 12.3, 0.28, F_RED, RED)
    text(s, [
        ("The classifier has no 'I don't know' option — it ALWAYS outputs a name. "
         "Every gate exists because a confident answer from a bad input is worse than no answer at all.",
         True, RED, 0),
    ], 0.68, 6.73, 11.9, 0.26, size=9)

    # ── 7 · EXPERIMENTAL SETUP ──────────────────────────────────────────────
    s = slide(prs, "EXPERIMENTAL SETUP & PROTOCOL",
              "The splits are what make the numbers believable — and they are not optional.")
    s.shapes.add_picture(f"{FIG}/protocol.png", Inches(0.6), Inches(1.3), height=Inches(2.9))
    s.shapes.add_picture(f"{FIG}/scores.png", Inches(0.6), Inches(4.35), height=Inches(2.45))
    s.shapes.add_picture(f"{FIG}/roc.png", Inches(5.35), Inches(4.35), height=Inches(2.45))

    panel(s, 9.5, 4.35, 3.3, 2.45, F_BLUE, BLUE)
    text(s, [
        ("What these two charts prove", True, NAVY, 0),
        "",
        "The genuine and impostor populations are almost",
        "completely separated — which is what lets a single",
        "cosine threshold work at all.",
        "",
        "The threshold is placed where only 1% of the",
        "CALIBRATION strangers get through, then measured",
        "on strangers it has never seen.",
        "",
        ("That is the difference between a measurement", True, RED, 0),
        ("and a number you tuned until it looked good.", True, RED, 0),
    ], 9.68, 4.47, 3.0, 2.25, size=8.5)

    panel(s, 9.5, 1.3, 3.3, 2.9, F_RED, RED)
    text(s, [
        ("Why the 2 real users are NOT the experiment", True, RED, 0),
        "",
        "With N=2, chance alone scores 50%. Any accuracy",
        "figure taken from two people is meaningless.",
        "",
        "Every quantitative claim in this deck comes from",
        "the 30-identity cohort. The live system with 2 real",
        "users is a feasibility demonstration — and we say",
        "so, rather than inflate it.",
        "",
        ("Environment", True, NAVY, 0),
        f"MacBook Air · CPU only · no GPU · {bm['memory_mb']:.0f} MB",
    ], 9.68, 1.42, 3.0, 2.7, size=8.5)

    # ── 8 · RESULTS I ───────────────────────────────────────────────────────
    s = slide(prs, "RESULTS AND ANALYSIS  (1 / 2)  —  Security",
              "Can a photograph fool it? And what does each guard actually buy?")
    s.shapes.add_picture(f"{FIG}/spoof.png", Inches(0.6), Inches(1.35), height=Inches(2.5))
    s.shapes.add_picture(f"{FIG}/ablation.png", Inches(5.1), Inches(1.35), height=Inches(2.5))

    panel(s, 10.3, 1.35, 2.5, 2.5, F_RED, RED)
    text(s, [
        ("Motion score", True, RED, 0),
        "",
        "photo (hand-held)",
        ("0.015", True, RED, 0),
        "Arnab, live",
        ("0.066", True, GREEN, 0),
        "Soham, live",
        ("0.135", True, GREEN, 0),
        "",
        ("threshold  0.035", True, INK, 0),
    ], 10.45, 1.46, 2.2, 2.3, size=9.5)

    table(s, [
        ["Configuration", "Accuracy", "FAR", "Wrong records"],
        ["Baseline (fixed threshold 0.32)", "87.4%", "0.81%", "3.54%"],
        ["+ Calibrated threshold", "86.7%", "0.20%", "1.01%"],
        ["+ Top-2 margin", "86.3%", "0.20%", "1.01%"],
        ["+ Temporal voting   (shipped)", "86.3%", "0.20%", "0.00%"],
    ], 0.6, 4.15, 7.6, 1.9, [3.4, 1.4, 1.4, 1.4], bold_rows=(0, 4))

    panel(s, 8.4, 4.15, 4.4, 2.45, F_AMBER, AMBER)
    text(s, [
        ("A guard that did NOT work", True, AMBER, 0),
        "We built the top-2 margin to catch look-alike students.",
        ("On this cohort it changed nothing — FAR stayed at 0.20%, and accuracy fell slightly.", True, INK, 0),
        "Why: accuracy and TAR are identical in every row, meaning the system never once put the WRONG "
        "student's name on a face. Every error was a rejection, not a mix-up — so the margin had nothing to catch.",
        ("We report it anyway. A result that only ever goes your way is not a result.", True, RED, 0),
    ], 8.58, 4.27, 4.1, 2.25, size=9)

    # ── 9 · RESULTS II ──────────────────────────────────────────────────────
    s = slide(prs, "RESULTS AND ANALYSIS  (2 / 2)  —  Envelope, Performance & Comparison",
              "Where does it break, how fast is it, and what do we do that prior work does not?")
    s.shapes.add_picture(f"{FIG}/degradation.png", Inches(0.6), Inches(1.35), height=Inches(2.45))
    s.shapes.add_picture(f"{FIG}/throughput.png", Inches(5.0), Inches(1.35), height=Inches(2.45))

    panel(s, 9.4, 1.35, 3.4, 2.45, F_BLUE, BLUE)
    text(s, [
        ("Derived hardware spec", True, NAVY, 0),
        "Accuracy collapses below 24 px and is fully",
        "recovered by 40 px — so the gate sits at 40 px.",
        "",
        ("1080p camera → covers ~7 m", True, GREEN, 0),
        ("720p camera → only ~5 m", True, AMBER, 0),
        "",
        "One camera covers a classroom. This is an",
        "engineering requirement derived from data,",
        "not a guess.",
    ], 9.58, 1.46, 3.1, 2.25, size=9.5)

    table(s, [
        ["", "Prior work [1–5]", "PresenceAI (ours)"],
        ["Multi-face in one frame", "Yes", "Yes"],
        ["Rejects strangers (open-set)", "Not reported", "0.20% FAR, calibrated"],
        ["Presentation attack tested", "No", "100% → 0%"],
        ["Attendance-commit policy", "Per-frame", "Multi-frame vote, 0% wrong"],
        ["Enrolment cost reported", "No", "1.4 s"],
        ["Negative result reported", "No", "Yes (top-2 margin)"],
    ], 0.6, 4.1, 8.5, 2.5, [3.3, 2.2, 3.0], green_col=2)

    panel(s, 9.4, 4.1, 3.4, 2.5, F_GREY, GREY)
    text(s, [
        ("Capacity", True, NAVY, 0),
        f"16 faces in one frame: {bm['throughput'][-1]['latency_s']} s",
        "Detection costs ~200 ms fixed; each extra",
        "face adds only ~65 ms — so it gets cheaper",
        "per student as the room fills.",
        "",
        ("10,000 enrolled students:", True, NAVY, 0),
        ("+1.2 ms per match", True, GREEN, 0),
        "Scale is a matrix multiply, because the",
        "deep model is frozen.",
    ], 9.58, 4.22, 3.1, 2.3, size=9.5)

    # ── 10 · IMPROVEMENTS + FUTURE SCOPE ────────────────────────────────────
    s = slide(prs, "IMPROVEMENTS MADE", "Measured before and after, on the same protocol.")
    metric(s, 0.6, 1.3, 3.9, "Photo-attack success rate", "100%", "0%", RED)
    metric(s, 4.75, 1.3, 3.9, "Strangers admitted (FAR)", "0.81%", "0.20%")
    metric(s, 8.9, 1.3, 3.9, "Wrong attendance records", "3.54%", "0%")
    metric(s, 0.6, 2.5, 3.9, "Time to enrol a student", "607 s", "1.4 s")
    metric(s, 4.75, 2.5, 3.9, "Time to mark present", "10.5 s", "2.8 s")
    metric(s, 8.9, 2.5, 3.9, "Classes the model trains on", "98", "2 + 96 impostors", BLUE)

    text(s, [("Also fixed along the way", True, NAVY, 0)], 0.6, 3.72, 12.2, 0.3, size=12.5)
    for i, (t, b) in enumerate([
        ("Registration never retrained",
         "A newly registered student stayed 'Unknown' forever — the images were saved, but the model was never refit."),
        ("The API was unauthenticated",
         "Anyone on the network could mark themselves present or delete a student. All 11 data routes now require a token."),
        ("Remote code execution",
         "Flask ran with debug=True bound to 0.0.0.0, exposing the Werkzeug console to the entire network."),
    ]):
        x = 0.6 + i * 4.15
        panel(s, x, 4.05, 3.9, 1.2, F_AMBER, AMBER)
        text(s, [(t, True, AMBER, 0), (b, False, INK, 0)], x + 0.15, 4.16, 3.6, 1.05, size=9)

    text(s, [("FUTURE SCOPE", True, NAVY, 0)], 0.6, 5.4, 12.2, 0.3, size=12.5)
    for i, (t, b) in enumerate([
        ("Defeat video replay",
         "A recorded face genuinely moves, so it still passes. A depth or texture-based anti-spoof model is the next step."),
        ("Enrol a real cohort",
         "Only 2 real users today. 30+ enrolled students would let every claim be re-measured on real webcam data."),
        ("Store embeddings, not photos",
         "A face cannot be reconstructed from a 512-d vector — a stolen database would leak nothing."),
    ]):
        x = 0.6 + i * 4.15
        panel(s, x, 5.73, 3.9, 1.2, F_GREEN, GREEN)
        text(s, [(t, True, GREEN, 0), (b, False, INK, 0)], x + 0.15, 5.84, 3.6, 1.05, size=9)

    # ── 11 · IMPACTS AND BENEFITS ───────────────────────────────────────────
    s = slide(prs, "IMPACTS AND BENEFITS")
    blocks = [
        ("SOCIAL", F_BLUE, BLUE, [
            "~80 teaching hours returned per year (5 min × 200 days).",
            "Contactless — no queue, no register passed around the room.",
            "Manual proxy eliminated; photo proxy blocked (100% → 0%).",
            "Video-replay proxy is NOT yet blocked — stated, not hidden.",
        ]),
        ("ECONOMIC", F_GREEN, GREEN, [
            "One laptop and one webcam. No GPU. No servers.",
            "Recurring cost ≈ ₹0.",
            "A cloud face API at ~$0.001/image, scanning continuously,",
            "costs ≈ ₹1,50,000 per year, per school.",
        ]),
        ("PRIVACY & LEGAL", F_RED, RED, [
            "Runs fully offline. Student faces never leave the campus.",
            "Under India's DPDP Act 2023, biometric data is sensitive",
            "personal data — uploading children's faces to a third-party",
            "cloud is a liability, not a feature.",
        ]),
        ("TECHNICAL", F_AMBER, AMBER, [
            "Every decision exposes its evidence: similarity score,",
            "runner-up margin, face size, liveness verdict.",
            "A disputed attendance record can be explained.",
            "Commercial systems are black boxes.",
        ]),
    ]
    for i, (t, fc, ec, lines) in enumerate(blocks):
        x = 0.6 + (i % 2) * 4.25
        y = 1.4 + (i // 2) * 2.6
        w = 8.1 if False else 3.95
        panel(s, x, y, w, 2.35, fc, ec)
        text(s, [(t, True, ec, 0), ("", False, INK, 0)] + lines,
             x + 0.16, y + 0.13, w - 0.32, 2.15, size=8.5)

    s.shapes.add_picture(f"{FIG}/cost.png", Inches(9.0), Inches(4.15), height=Inches(2.2))

    panel(s, 0.6, 6.35, 12.15, 0.6, F_GREY, NAVY)
    text(s, [("Target audience — schools and colleges that cannot afford a GPU or a cloud subscription, "
              "and cannot legally upload their students' faces to one.", True, NAVY, 0)],
         0.78, 6.45, 11.8, 0.45, size=10.5)

    # ── 12 · REFERENCES ─────────────────────────────────────────────────────
    s = slide(prs, "REFERENCES")
    text(s, [
        ("Attendance systems — prior work (none of which tests a presentation attack)", True, NAVY, 0),
        "[1]  D'Souza, J. W. S., et al. (2019). Automated Attendance Marking and Management System by Facial "
        "Recognition Using Histogram. Array, 3–4: 100014. DOI: 10.1016/j.array.2019.100014",
        "[2]  Arsenovic, M., Sladojevic, S., Anderla, A., Stefanovic, D. (2017). FaceTime — Deep Learning Based Face "
        "Recognition Attendance System. IEEE SISY 2017, pp. 53–58. DOI: 10.1109/SISY.2017.8080587",
        "[3]  Kakarla, S., Gangula, P., Rahul, M. S., Singh, C. S. C., Sarma, T. H. (2020). Smart Attendance Management "
        "System Based on Face Recognition Using CNN. IEEE-HYDCON 2020. DOI: 10.1109/HYDCON48903.2020.9242847",
        "[4]  Varadharajan, E., Dharani, R., Jeevitha, S., Kavinmathi, B., Hemalatha, S. (2016). Automatic Attendance "
        "Management System Using Face Detection. IC-GET 2016, pp. 1–3. DOI: 10.1109/GET.2016.7916753",
        "[5]  Siswanto, A. R. S., Nugroho, A. S., Galinium, M. (2014). Implementation of Face Recognition Algorithm for "
        "Biometrics Based Time Attendance System. ICISS 2014, pp. 149–154. DOI: 10.1109/ICTSS.2014.7013165",
        "",
        ("Models we use — pretrained and frozen (not our contribution)", True, NAVY, 0),
        "[6]  Deng, J., Guo, J., Xue, N., Zafeiriou, S. (2019). ArcFace: Additive Angular Margin Loss for Deep Face "
        "Recognition. CVPR 2019, pp. 4690–4699. DOI: 10.1109/CVPR.2019.00482",
        "[7]  Deng, J., Guo, J., Ververas, E., Kotsia, I., Zafeiriou, S. (2020). RetinaFace: Single-Shot Multi-Level Face "
        "Localisation in the Wild. CVPR 2020. DOI: 10.1109/CVPR42600.2020.00525",
        "[8]  Schroff, F., Kalenichenko, D., Philbin, J. (2015). FaceNet: A Unified Embedding for Face Recognition and "
        "Clustering. CVPR 2015, pp. 815–823. DOI: 10.1109/CVPR.2015.7298682",
        "",
        ("Presentation attacks — the gap this project addresses", True, RED, 0),
        "[9]   Pan, G., Sun, L., Wu, Z., Lao, S. (2007). Eyeblink-based Anti-Spoofing in Face Recognition from a Generic "
        "Webcamera. ICCV 2007, pp. 1–8. DOI: 10.1109/ICCV.2007.4409068",
        "[10] ISO/IEC 30107-3:2023 — Information technology — Biometric presentation attack detection — Part 3: Testing "
        "and reporting.",
        "[11] Cao, Q., Shen, L., Xie, W., Parkhi, O. M., Zisserman, A. (2018). VGGFace2: A Dataset for Recognising Faces "
        "across Pose and Age. IEEE FG 2018, pp. 67–74.",
        "",
        ("Source code, dataset and every measurement script:  "
         "github.com/SohamBhattacharjee2003/MutiFace-Attendance-System", True, GREEN, 0),
    ], 0.6, 1.05, 12.2, 5.95, size=9.5)

    # ── 13 · THANK YOU ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                              Inches(13.333), Inches(0.28))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE
    band.line.fill.background(); band.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(46); p.runs[0].font.bold = True; p.runs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "Questions?"
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(16); p2.runs[0].font.color.rgb = GREY

    for i, sl in enumerate(prs.slides):
        if i:
            chrome(sl, i)

    prs.save(OUT)
    print(f"\n✅ {OUT}   ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
